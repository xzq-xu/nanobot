"""Document text extraction utilities for nanobot."""

from pathlib import Path

from loguru import logger

try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None  # type: ignore

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None  # type: ignore

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None  # type: ignore

try:
    from pptx import Presentation as PptxPresentation
except ImportError:
    PptxPresentation = None  # type: ignore


# Supported file extensions for text extraction
SUPPORTED_EXTENSIONS: set[str] = {
    # Document formats
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
    # Text formats
    ".txt",
    ".md",
    ".csv",
    ".json",
    ".xml",
    ".html",
    ".htm",
    ".log",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
    # Image formats (for future OCR support)
    ".png",
    ".jpg",
    ".jpeg",
    ".gif",
    ".webp",
}

_MAX_TEXT_LENGTH = 200_000


def extract_text(path: Path) -> str | None:
    """Extract text from a file.

    Args:
        path: Path to the file.

    Returns:
        Extracted text as string, None for unsupported types,
        or error string for failures.
    """
    if not isinstance(path, Path):
        path = Path(path)

    if not path.exists():
        return f"[error: file not found: {path}]"

    ext = path.suffix.lower()

    # Document formats
    if ext == ".pdf":
        if PdfReader is None:
            return "[error: pypdf not installed]"
        return _extract_pdf(path)
    elif ext == ".docx":
        if DocxDocument is None:
            return "[error: python-docx not installed]"
        return _extract_docx(path)
    elif ext == ".xlsx":
        if load_workbook is None:
            return "[error: openpyxl not installed]"
        return _extract_xlsx(path)
    elif ext == ".pptx":
        if PptxPresentation is None:
            return "[error: python-pptx not installed]"
        return _extract_pptx(path)
    elif _is_text_extension(ext):
        return _extract_text_file(path)
    elif ext in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
        # Image files - for future OCR support
        return f"[image: {path.name}]"
    else:
        # Unsupported extension
        return None


def _extract_pdf(path: Path) -> str:
    """Extract text from PDF using pypdf."""
    try:
        reader = PdfReader(path)
        pages: list[str] = []
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            pages.append(f"--- Page {i} ---\n{text}")
        return _truncate("\n\n".join(pages), _MAX_TEXT_LENGTH)
    except Exception as e:
        logger.error("Failed to extract PDF {}: {}", path, e)
        return f"[error: failed to extract PDF: {e!s}]"


def _extract_docx(path: Path) -> str:
    """Extract text from DOCX using python-docx."""
    try:
        doc = DocxDocument(path)
        paragraphs: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
        return _truncate("\n\n".join(paragraphs), _MAX_TEXT_LENGTH)
    except Exception as e:
        logger.error("Failed to extract DOCX {}: {}", path, e)
        return f"[error: failed to extract DOCX: {e!s}]"


def _extract_xlsx(path: Path) -> str:
    """Extract text from XLSX using openpyxl."""
    try:
        wb = load_workbook(path, read_only=True, data_only=True)
        sheets: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                sheets.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
        wb.close()
        return _truncate("\n\n".join(sheets), _MAX_TEXT_LENGTH)
    except Exception as e:
        logger.error("Failed to extract XLSX {}: {}", path, e)
        return f"[error: failed to extract XLSX: {e!s}]"


def _extract_pptx(path: Path) -> str:
    """Extract text from PPTX using python-pptx."""
    try:
        prs = PptxPresentation(path)
        slides: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            if slide_text:
                slides.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))
        return _truncate("\n\n".join(slides), _MAX_TEXT_LENGTH)
    except Exception as e:
        logger.error("Failed to extract PPTX {}: {}", path, e)
        return f"[error: failed to extract PPTX: {e!s}]"


def _extract_text_file(path: Path) -> str:
    """Extract text from a plain text file."""
    try:
        # Try UTF-8 first, then latin-1 fallback
        try:
            content = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            content = path.read_text(encoding="latin-1")
        return _truncate(content, _MAX_TEXT_LENGTH)
    except Exception as e:
        logger.error("Failed to read text file {}: {}", path, e)
        return f"[error: failed to read file: {e!s}]"


def _truncate(text: str, max_length: int) -> str:
    """Truncate text with a suffix indicating truncation."""
    if len(text) <= max_length:
        return text
    return text[:max_length] + f"... (truncated, {len(text)} chars total)"


def _is_text_extension(ext: str) -> bool:
    """Check if extension is a text format."""
    return ext in {
        ".txt",
        ".md",
        ".csv",
        ".json",
        ".xml",
        ".html",
        ".htm",
        ".log",
        ".yaml",
        ".yml",
        ".toml",
        ".ini",
        ".cfg",
    }
